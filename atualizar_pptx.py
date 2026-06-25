"""
Atualiza analise_atendimentos_2026.pptx com dados de 26/05/2026.
Dados novos vs antigos (20/05):
  Vendas Mai: 4384 → 5211
  F.Produção Mai: 34 → 44  | Total: 674 → 686
  F.Transporte Mai: 104 → 139 | Total: 599 → 631
  F.Compras Mai: 46 → 65  | Total: 283 → 305
"""
import copy, re
from pptx import Presentation
from pptx.chart.data import ChartData

SRC  = 'C:/Users/DELL/Downloads/analise_atendimentos_2026.pptx'
DEST = 'C:/Users/DELL/Downloads/analise_atendimentos_2026_26maio.pptx'

prs = Presentation(SRC)

# ── helpers ────────────────────────────────────────────────────────────────
def replace_in_shape(shape, old, new):
    """Substitui texto em todos os runs de um shape, preservando formatação."""
    if not shape.has_text_frame:
        return False
    found = False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)
                found = True
    return found

def replace_in_slide(slide, old, new):
    for shape in slide.shapes:
        replace_in_shape(shape, old, new)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            if old in run.text:
                                run.text = run.text.replace(old, new)

def get_chart(slide):
    for shape in slide.shapes:
        if shape.has_chart:
            return shape.chart
    return None

def update_column_chart(slide, series_name, values):
    chart = get_chart(slide)
    if not chart:
        return
    cd = ChartData()
    cd.categories = ['Jan', 'Fev', 'Mar', 'Abr', 'Mai']
    cd.add_series(series_name, values)
    chart.replace_data(cd)

def update_line_chart(slide, series_data):
    """series_data = [(name, [vals]), ...]"""
    chart = get_chart(slide)
    if not chart:
        return
    cd = ChartData()
    cd.categories = ['Jan', 'Fev', 'Mar', 'Abr', 'Mai']
    for name, vals in series_data:
        cd.add_series(name, vals)
    chart.replace_data(cd)

slides = prs.slides

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Título
# ══════════════════════════════════════════════════════════════════════════
s1 = slides[0]
replace_in_slide(s1, '20/05/2026', '26/05/2026')

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Vendas
# ══════════════════════════════════════════════════════════════════════════
s2 = slides[1]
update_column_chart(s2, 'Pedidos', (8152, 8015, 8712, 7906, 5211))
replace_in_slide(s2, '37.743', '37.996')
replace_in_slide(s2, '8.340',  '8.196')
replace_in_slide(s2, '8.839',  '8.712')
replace_in_slide(s2, 'até 20/05/2026', 'até 26/05/2026')

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Falha Produção
# ══════════════════════════════════════════════════════════════════════════
s3 = slides[2]
update_column_chart(s3, 'Atendimentos', (108, 164, 271, 99, 44))
replace_in_slide(s3, '−87,4%', '−83,8%')
replace_in_slide(s3, 'Mar: 3,05%', 'Mar: 3,11%')
replace_in_slide(s3, 'Mai: 0,78%', 'Mai: 0,84%')
replace_in_slide(s3, 'Total no ano: 674 atendimentos', 'Total no ano: 686 atendimentos')

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Falha Transporte
# ══════════════════════════════════════════════════════════════════════════
s4 = slides[3]
update_column_chart(s4, 'Atendimentos', (111, 84, 160, 137, 139))
replace_in_slide(s4, '−35,4%', '−13,1%')
replace_in_slide(s4, 'Mar: 1,82%', 'Mar: 1,84%')
replace_in_slide(s4, 'Mai: 2,37%', 'Mai: 2,67%')
replace_in_slide(s4, 'Total no ano: 599 atendimentos', 'Total no ano: 631 atendimentos')
replace_in_slide(s4, 'Volume caiu mas a taxa sobre vendas aumentou. Maio é parcial — acompanhar o fechamento.',
                     'Volume estável em relação a Abril (137→139). Taxa subiu para 2,67% — categoria requer atenção no fechamento de Maio.')

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Falha Compras
# ══════════════════════════════════════════════════════════════════════════
s5 = slides[4]
update_column_chart(s5, 'Atendimentos', (56, 49, 67, 68, 65))
replace_in_slide(s5, 'Jan 0,68%', 'Jan 0,69%')
replace_in_slide(s5, 'Fev 0,60%', 'Fev 0,61%')
replace_in_slide(s5, 'Mar 0,74%', 'Mar 0,77%')
replace_in_slide(s5, 'Abr 0,83%', 'Abr 0,86%')
replace_in_slide(s5, 'Mai 1,05%', 'Mai 1,25%')
replace_in_slide(s5, 'há 4 meses', 'há 3 meses')
replace_in_slide(s5, 'Total no ano: 283 atendimentos', 'Total no ano: 305 atendimentos')
replace_in_slide(s5, 'Menor das 3 em volume, mas a única em ALTA proporcional contínua. Investigar causa antes que vire pico.',
                     'Taxa atingiu 1,25% em Maio — maior do ano. Única categoria em alta contínua. Investigar causa antes que vire pico.')
replace_in_slide(s5, 'Categoria em alta proporcional há 4 meses consecutivos',
                     'Categoria em alta proporcional — taxa recorde em Maio')

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Comparativo taxa
# ══════════════════════════════════════════════════════════════════════════
s6 = slides[5]
update_line_chart(s6, [
    ('Falha Produção',   [1.32, 2.05, 3.11, 1.25, 0.84]),
    ('Falha Transporte', [1.36, 1.05, 1.84, 1.73, 2.67]),
    ('Falha Compras',    [0.69, 0.61, 0.77, 0.86, 1.25]),
])
replace_in_slide(s6, 'Produção em queda · Transporte e Compras em alta proporcional (Maio parcial)',
                     'Produção em queda · Compras em recorde · Transporte acima de 2,6% (Maio até 26/05)')

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Conclusões
# ══════════════════════════════════════════════════════════════════════════
s10 = slides[9]
replace_in_slide(s10, 'Redução de 74% na taxa proporcional (Mar→Mai). Equipe entregou resultado real.',
                      'Redução de 73% na taxa proporcional (Mar→Mai). Equipe entregou resultado real.')
replace_in_slide(s10, 'Volume caindo, mas taxa proporcional ainda não evoluiu. Acompanhar.',
                      'Volume estável, taxa em 2,67% — acima dos meses anteriores. Acompanhar.')
replace_in_slide(s10, 'Única em alta proporcional contínua há 4 meses. Investigar causa-raiz.',
                      'Taxa em 1,25% — recorde do ano. Alta contínua há 3 meses. Investigar causa-raiz.')
replace_in_slide(s10, 'Jan–Mai 2026 · Reclame Aqui · Compilado em 20/05/2026',
                      'Jan–Mai 2026 · Reclame Aqui · Compilado em 26/05/2026')

# ── Salva ──────────────────────────────────────────────────────────────────
prs.save(DEST)
print(f'Salvo: {DEST}')

# ── Validação básica ───────────────────────────────────────────────────────
prs2 = Presentation(DEST)
checks = [
    (1, '26/05/2026'),
    (1, '37.996'),
    (1, '8.196'),
    (2, '686'),
    (3, '631'),
    (4, '305'),
    (4, '1,25%'),
    (5, '26/05/2026'),
]
print('\nValidação:')
for slide_idx, term in checks:
    slide = prs2.slides[slide_idx]
    all_text = ' '.join(
        run.text
        for shape in slide.shapes if shape.has_text_frame
        for para in shape.text_frame.paragraphs
        for run in para.runs
    )
    ok = '✅' if term in all_text else '❌'
    print(f'  {ok} Slide {slide_idx+1}: "{term}"')
