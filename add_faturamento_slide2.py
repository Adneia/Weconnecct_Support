"""
Adiciona faturamento ao slide 2: shrink chart, adiciona 3 KPI cards de receita abaixo,
e reorganiza os KPI cards direitos para caber 5 métricas.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

SRC  = 'C:/Users/DELL/Downloads/analise_atendimentos_2026_26maio.pptx'
DEST = 'C:/Users/DELL/Downloads/analise_atendimentos_2026_26maio_v2.pptx'

def emu(inch): return int(inch * 914400)

prs = Presentation(SRC)
s2  = prs.slides[1]

# ── 1. Encolhe chart de H=4.80 para H=3.80 ────────────────────────────────
for sh in s2.shapes:
    if sh.has_chart:
        sh.height = emu(3.80)
        break

# ── 2. Atualiza subtítulo ──────────────────────────────────────────────────
for sh in s2.shapes:
    if sh.has_text_frame and sh.name == 'Text 1':
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                if 'Volume de pedidos' in run.text:
                    run.text = 'Volume de pedidos faturados e receita por mês'

# ── 3. Adiciona 3 KPI cards de faturamento abaixo do chart ─────────────────
# Área do chart: L=0.70, T=1.80, W=8.00 → bottom agora = 1.80+3.80 = 5.60
# 3 cards de L=0.70 até L=8.70, cada W=2.55, gap=0.07, T=5.65, H=1.10

fat_cards = [
    ('FATURAMENTO TOTAL',  'R$ 7.642.344', 'Jan–Mai 2026'),
    ('TICKET MÉDIO',       'R$ 201',       'por pedido'),
    ('PICO · MARÇO',       'R$ 1.745.968', 'maior receita do ano'),
]

card_l_starts = [0.70, 3.32, 5.94]
card_t  = 5.65
card_w  = 2.55
card_h  = 1.10
fill_hex = '1A3A5C'   # azul escuro igual ao tema

for i, (label, value, sub) in enumerate(fat_cards):
    l = card_l_starts[i]

    # fundo
    bg = s2.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        emu(l), emu(card_t), emu(card_w), emu(card_h)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x1A, 0x3A, 0x5C)
    bg.line.fill.background()

    # label (pequeno, topo)
    tb_label = s2.shapes.add_textbox(
        emu(l + 0.08), emu(card_t + 0.06), emu(card_w - 0.16), emu(0.22)
    )
    p = tb_label.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = label
    run.font.bold  = True
    run.font.size  = Pt(7)
    run.font.color.rgb = RGBColor(0xA8, 0xC8, 0xE8)
    run.font.name  = 'Arial'

    # valor (grande, meio)
    tb_val = s2.shapes.add_textbox(
        emu(l + 0.08), emu(card_t + 0.28), emu(card_w - 0.16), emu(0.50)
    )
    p = tb_val.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = value
    run.font.bold  = True
    run.font.size  = Pt(18)
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.name  = 'Arial'

    # sub (pequeno, base)
    tb_sub = s2.shapes.add_textbox(
        emu(l + 0.08), emu(card_t + 0.80), emu(card_w - 0.16), emu(0.22)
    )
    p = tb_sub.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = sub
    run.font.size  = Pt(7)
    run.font.color.rgb = RGBColor(0xA8, 0xC8, 0xE8)
    run.font.name  = 'Arial'

# ── 4. Reposiciona os 3 KPI cards direitos para caber em T=1.80–5.55 ───────
# Original: T= 1.90, 3.40, 4.90 com H=1.30
# Novo:     T= 1.80, 3.05, 4.30 com H=1.15  (gap 0.10)
right_card_shapes = {
    'Shape 2': (1.80, 1.15),
    'Shape 6': (3.05, 1.15),
    'Shape 10': (4.30, 1.15),
}
right_text_tops = {
    'Text 3': 1.88,  'Text 4': 2.16,  'Text 5': 2.68,   # card 1
    'Text 7': 3.13,  'Text 8': 3.41,  'Text 9': 3.93,   # card 2
    'Text 11': 4.38, 'Text 12': 4.66, 'Text 13': 5.18,  # card 3
}

for sh in s2.shapes:
    if sh.name in right_card_shapes:
        new_t, new_h = right_card_shapes[sh.name]
        sh.top    = emu(new_t)
        sh.height = emu(new_h)
    if sh.name in right_text_tops:
        sh.top = emu(right_text_tops[sh.name])
        # reduz font do número grande (Text 4, 8, 12) se necessário
        if sh.name in ('Text 4', 'Text 8', 'Text 12') and sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size and run.font.size > Pt(20):
                        run.font.size = Pt(22)

prs.save(DEST)
print(f'Salvo: {DEST}')
